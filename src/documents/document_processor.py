"""
Document Processing Module
Extracts text from PDF, Word (DOCX), and PowerPoint (PPTX) files
Supports repair manuals, bulletins, and technical documentation

Phase 1-2 Enhancement: Integrated semantic chunking for better RAG retrieval
"""
import re
from pathlib import Path
from typing import List, Dict, Optional, Tuple

# PyPDF2 may be available as 'PyPDF2' or as the newer 'pypdf' package; try both for compatibility.
try:
    import PyPDF2  # type: ignore
    PYPDF2_AVAILABLE = True
except ImportError:
    try:
        import pypdf as PyPDF2  # type: ignore
        PYPDF2_AVAILABLE = True
    except ImportError:
        PyPDF2 = None  # type: ignore
        PYPDF2_AVAILABLE = False

# PDF text extraction with pdfplumber
try:
    import pdfplumber  # type: ignore
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    pdfplumber = None  # type: ignore
    PDFPLUMBER_AVAILABLE = False

from src.utils.logger import setup_logger

# Word document support
try:
    from docx import Document as DocxDocument  # type: ignore
    DOCX_AVAILABLE = True
except ImportError:
    DocxDocument = None  # type: ignore
    DOCX_AVAILABLE = False

# PowerPoint support
try:
    from pptx import Presentation  # type: ignore
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None  # type: ignore
    PPTX_AVAILABLE = False

# Excel support
try:
    import openpyxl  # type: ignore # pyright: ignore[import]
    EXCEL_AVAILABLE = True
except ImportError:
    openpyxl = None  # type: ignore
    EXCEL_AVAILABLE = False

# Semantic chunking for Phase 1-2 enhancement
from src.documents.semantic_chunker import SemanticChunker, DocumentType, ChunkMetadata
from src.documents.product_extractor import ProductExtractor

logger = setup_logger(__name__)

# Supported file extensions
SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls'}


class DocumentProcessor:
    """Process documents (PDF, Word, PowerPoint) and extract text"""
    
    def __init__(self):
        """Initialize document processor"""
        logger.info(f"DocumentProcessor initialized")
        logger.info(f"  - PDF support: {PYPDF2_AVAILABLE or PDFPLUMBER_AVAILABLE}")
        logger.info(f"  - DOCX support: {DOCX_AVAILABLE}")
        logger.info(f"  - PPTX support: {PPTX_AVAILABLE}")
        logger.info(f"  - EXCEL support: {EXCEL_AVAILABLE}")
        
        # Initialize semantic chunker for Phase 1-2 enhancement
        self.semantic_chunker = SemanticChunker(
            chunk_size=400,
            chunk_overlap=100,
            max_recursion_depth=3
        )
        
        # Initialize product extractor for Metadata Enrichment
        self.product_extractor = ProductExtractor()
        
        logger.info(f"  - Semantic chunking: ENABLED (Phase 1-2 enhancement)")
        logger.info(f"  - Product recognition: ENABLED")
    
    # =========================================================================
    # PDF PROCESSING
    # =========================================================================
    
    def extract_text_pypdf2(self, pdf_path: Path) -> str:
        """
        Extract text using PyPDF2 (faster, basic)
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            Extracted text
        """
        try:
            text_content = []
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(text)
                        
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error extracting text from {pdf_path}: {e}")
            return ""
    
    def extract_text_pdfplumber(self, pdf_path: Path, extract_tables: bool = True) -> Dict:
        """
        Extract text and tables using pdfplumber (better quality)
        
        Args:
            pdf_path: Path to PDF file
            extract_tables: Whether to extract tables
            
        Returns:
            Dictionary with text and tables
        """
        if not PDFPLUMBER_AVAILABLE:
            logger.warning("pdfplumber not available, falling back to PyPDF2")
            return {"text": self.extract_text_pypdf2(pdf_path), "tables": "", "combined": self.extract_text_pypdf2(pdf_path)}
            
        try:
            text_content = []
            tables_content = []
            
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    # Extract text
                    text = page.extract_text()
                    if text and text.strip():
                        text_content.append(f"--- Page {page_num + 1} ---\n{text}")
                    
                    # Extract tables if requested
                    if extract_tables:
                        tables = page.extract_tables()
                        for table_idx, table in enumerate(tables):
                            if table:
                                table_text = self._format_table(table)
                                tables_content.append(
                                    f"--- Page {page_num + 1}, Table {table_idx + 1} ---\n{table_text}"
                                )
            
            return {
                "text": "\n\n".join(text_content),
                "tables": "\n\n".join(tables_content) if tables_content else "",
                "combined": "\n\n".join(text_content + tables_content)
            }
        except Exception as e:
            logger.error(f"Error extracting from {pdf_path}: {e}")
            return {"text": "", "tables": "", "combined": ""}
    
    def _format_table(self, table: List[List[str]]) -> str:
        """Format extracted table as text"""
        if not table:
            return ""
        
        # Convert table to text representation
        formatted_rows = []
        for row in table:
            # Clean None values
            cleaned_row = [str(cell).strip() if cell else "" for cell in row]
            formatted_rows.append(" | ".join(cleaned_row))
        
        return "\n".join(formatted_rows)
    
    def process_pdf(
        self,
        pdf_path: Path,
        method: str = "pdfplumber",
        extract_tables: bool = True
    ) -> Optional[str]:
        """
        Process a PDF document and extract text
        
        Args:
            pdf_path: Path to PDF file
            method: Extraction method ('pypdf2' or 'pdfplumber')
            extract_tables: Whether to extract tables
            
        Returns:
            Extracted text content
        """
        logger.info(f"Processing PDF: {pdf_path.name}")
        
        if method == "pdfplumber":
            extracted = self.extract_text_pdfplumber(pdf_path, extract_tables)
            return extracted["combined"]
        else:
            return self.extract_text_pypdf2(pdf_path)
    
    # =========================================================================
    # WORD (DOCX) PROCESSING
    # =========================================================================
    
    def process_docx(self, docx_path: Path) -> Optional[str]:
        """
        Extract text from Word document (DOCX)
        
        Args:
            docx_path: Path to DOCX file
            
        Returns:
            Extracted text content
        """
        if not DOCX_AVAILABLE:
            logger.error("python-docx not installed. Cannot process Word files.")
            return None
        
        logger.info(f"Processing Word: {docx_path.name}")
        
        try:
            doc = DocxDocument(docx_path)
            text_content = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_text = []
                table_text.append(f"\n--- Table {table_idx + 1} ---")
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                text_content.append("\n".join(table_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"Error extracting from {docx_path}: {e}")
            return None
    
    # =========================================================================
    # POWERPOINT (PPTX) PROCESSING
    # =========================================================================
    
    def process_pptx(self, pptx_path: Path) -> Optional[str]:
        """
        Extract text from PowerPoint presentation (PPTX)
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Extracted text content
        """
        if not PPTX_AVAILABLE:
            logger.error("python-pptx not installed. Cannot process PowerPoint files.")
            return None
        
        logger.info(f"Processing PowerPoint: {pptx_path.name}")
        
        try:
            prs = Presentation(pptx_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    # Extract text from text frames
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        table_text = []
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            table_text.append(" | ".join(row_text))
                        if table_text:
                            slide_text.append("\n".join(table_text))
                
                if len(slide_text) > 1:  # More than just slide header
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"Error extracting from {pptx_path}: {e}")
            return None
    
    # =========================================================================
    # EXCEL (XLSX/XLS) PROCESSING
    # =========================================================================
    
    def extract_text_excel(self, excel_path: Path) -> Optional[str]:
        """
        Extract text from Excel files (.xlsx, .xls)
        Reads all sheets and converts to text format
        
        Args:
            excel_path: Path to Excel file
            
        Returns:
            Extracted text from all sheets
        """
        if not EXCEL_AVAILABLE:
            logger.error("openpyxl not installed. Cannot process Excel files.")
            return None
        
        logger.info(f"Processing Excel: {excel_path.name}")
        
        try:
            text_content = []
            
            if excel_path.suffix.lower() == '.xlsx':
                from openpyxl import load_workbook  # type: ignore
                workbook = load_workbook(excel_path, data_only=True)
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_content.append(f"\n=== Sheet: {sheet_name} ===")
                    
                    for row in sheet.iter_rows(values_only=True):
                        # Filter out None values and convert to string
                        row_text = " | ".join(
                            str(cell) for cell in row if cell is not None
                        )
                        if row_text.strip():
                            text_content.append(row_text)
            
            elif excel_path.suffix.lower() == '.xls':
                # For .xls files, try with openpyxl first, fallback to basic reading
                try:
                    from openpyxl import load_workbook  # type: ignore
                    workbook = load_workbook(excel_path, data_only=True)
                    
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        text_content.append(f"\n=== Sheet: {sheet_name} ===")
                        
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " | ".join(
                                str(cell) for cell in row if cell is not None
                            )
                            if row_text.strip():
                                text_content.append(row_text)
                except:
                    logger.warning(f"Could not read .xls file {excel_path}, returning empty")
                    return ""
            
            return "\n".join(text_content)
        
        except Exception as e:
            logger.error(f"Error extracting text from Excel {excel_path}: {e}")
            return None
    
    # =========================================================================
    # UNIFIED DOCUMENT PROCESSING
    # =========================================================================
    
    def process_document(
        self,
        file_path: Path,
        extract_tables: bool = True,
        enable_semantic_chunking: bool = True
    ) -> Optional[Dict]:
        """
        Process any supported document type
        
        Args:
            file_path: Path to document file
            extract_tables: Whether to extract tables (PDF only)
            enable_semantic_chunking: Whether to apply semantic chunking (Phase 1-2)
            
        Returns:
            Dictionary with document metadata and content
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return None
        
        ext = file_path.suffix.lower()
        
        # Extract text based on file type
        if ext == '.pdf':
            text = self.process_pdf(file_path, extract_tables=extract_tables)
        elif ext in ['.docx', '.doc']:
            text = self.process_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            text = self.process_pptx(file_path)
        elif ext in ['.xlsx', '.xls']:
            text = self.extract_text_excel(file_path)
        else:
            logger.error(f"Unsupported file type: {ext}")
            return None
        
        if not text or len(text.strip()) < 50:
            logger.warning(f"No meaningful content extracted from {file_path.name}")
            return None
        
        # Clean text
        text = self.clean_text(text)
        
        # Extract metadata
        metadata = self.extract_metadata(file_path, text)
        
        # NEW: Get full product metadata from intelligent extractor
        product_metadata = self.product_extractor.get_product_metadata(
            filename=file_path.name,
            content=text[:2000]
        )
        
        # Log product detection
        logger.info(f"  Product detected: {product_metadata['product_family']} "
                    f"(confidence: {product_metadata['confidence']:.2f}, "
                    f"generic: {product_metadata['is_generic']})")
        
        # Phase 1-2: Apply semantic chunking if enabled
        chunks = None
        chunk_count = 0
        if enable_semantic_chunking:
            try:
                # Detect document type
                doc_type = self._map_to_document_type(metadata.get("type", "unknown"))
                
                # Apply semantic chunking with full product metadata
                chunks = self.semantic_chunker.chunk_document(
                    text=text,
                    source_filename=file_path.name,
                    doc_type=doc_type,
                    product_metadata=product_metadata  # NEW: Pass full metadata
                )
                chunk_count = len(chunks)
                logger.info(f"  Semantic chunking: Created {chunk_count} chunks")
            except Exception as e:
                logger.error(f"Error in semantic chunking: {e}")
                chunks = None
        
        return {
            "filename": file_path.name,
            "filepath": str(file_path),
            "file_type": ext.replace('.', ''),
            "text": text,
            "metadata": metadata,
            "product_metadata": product_metadata,  # NEW: Include for reference
            "word_count": len(text.split()),
            "char_count": len(text),
            "chunks": chunks,  # Phase 1-2: Semantic chunks
            "chunk_count": chunk_count
        }
    
    def clean_text(self, text: str) -> str:
        """
        Clean and normalize extracted text while preserving structure
        """
        # 1. Normalize line endings
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        
        # 2. Preserve paragraph breaks (double newlines)
        # Split by double newlines to get paragraphs
        paragraphs = re.split(r'\n\s*\n', text)
        
        cleaned_paragraphs = []
        for p in paragraphs:
            if not p.strip():
                continue
            # For each paragraph, replace single newlines with spaces (unwrap text)
            # but keep the paragraph itself as a unit
            # Normalize internal whitespace
            cleaned_p = re.sub(r'\s+', ' ', p).strip()
            cleaned_paragraphs.append(cleaned_p)
            
        # 3. Rejoin with double newlines
        return '\n\n'.join(cleaned_paragraphs)
    
    def extract_metadata(self, file_path: Path, text: str) -> Dict:
        """
        Extract metadata from file and text
        
        Args:
            file_path: Path to file
            text: Extracted text
            
        Returns:
            Metadata dictionary
        """
        metadata = {
            "source": file_path.name,
            "type": self._infer_document_type(file_path),
            "format": file_path.suffix.lower().replace('.', ''),
            # Metadata Enrichment: Smart Product Recognition
            "product_categories": self.product_extractor.get_product_categories(file_path.name, text[:1000])
        }
        
        # Try to extract product references from filename and text
        part_numbers = self._extract_part_numbers(file_path.name + " " + text[:1000])
        if part_numbers:
            metadata["related_products"] = part_numbers
        
        # Try to extract language
        metadata["language"] = self._detect_language(text[:500])
        
        return metadata
    
    def _infer_document_type(self, file_path: Path) -> str:
        """Infer document type from path"""
        path_str = str(file_path).lower()
        
        if "manual" in path_str:
            return "manual"
        elif "bulletin" in path_str:
            return "bulletin"
        elif "guide" in path_str:
            return "guide"
        elif "training" in path_str or "presentation" in path_str:
            return "training"
        else:
            return "unknown"
    
    def _map_to_document_type(self, doc_type_str: str) -> DocumentType:
        """Map document type string to DocumentType enum"""
        doc_type_lower = doc_type_str.lower()
        
        if "bulletin" in doc_type_lower:
            return DocumentType.SERVICE_BULLETIN
        elif "guide" in doc_type_lower or "troubleshoot" in doc_type_lower:
            return DocumentType.TROUBLESHOOTING_GUIDE
        elif "catalog" in doc_type_lower or "parts" in doc_type_lower:
            return DocumentType.PARTS_CATALOG
        elif "safety" in doc_type_lower:
            return DocumentType.SAFETY_DOCUMENT
        else:
            return DocumentType.TECHNICAL_MANUAL  # Default
    
    def _extract_part_numbers(self, text: str) -> List[str]:
        """Extract Desoutter part numbers from text"""
        patterns = [
            r'\b(6151[0-9]{6})\b',
            r'\b(8920[0-9]{6})\b',
            r'\b(6153[0-9]{6})\b',
            r'\b(6154[0-9]{6})\b',
        ]
        
        found = set()
        for pattern in patterns:
            matches = re.findall(pattern, text)
            found.update(matches)
        
        return sorted(list(found))
    
    def _detect_language(self, text: str) -> str:
        """Simple language detection"""
        # Turkish specific characters
        turkish_chars = ['ş', 'Ş', 'ğ', 'Ğ', 'ı', 'İ', 'ö', 'Ö', 'ü', 'Ü', 'ç', 'Ç']
        
        # Count Turkish characters
        turkish_count = sum(1 for char in text if char in turkish_chars)
        
        if turkish_count > 5:
            return "tr"
        else:
            return "en"
    
    def process_directory(
        self,
        directory: Path,
        extract_tables: bool = True,
        enable_semantic_chunking: bool = True
    ) -> List[Dict]:
        """
        Process all supported documents in a directory
        
        Args:
            directory: Directory containing documents
            extract_tables: Whether to extract tables
            enable_semantic_chunking: Whether to apply semantic chunking (Phase 1-2)
            
        Returns:
            List of processed documents
        """
        logger.info(f"Processing documents in: {directory}")
        
        # Find all supported files
        all_files = []
        for ext in SUPPORTED_EXTENSIONS:
            all_files.extend(directory.glob(f"*{ext}"))
            all_files.extend(directory.glob(f"*{ext.upper()}"))
        
        logger.info(f"Found {len(all_files)} supported documents")
        
        documents = []
        for doc_file in all_files:
            doc = self.process_document(
                doc_file, 
                extract_tables=extract_tables,
                enable_semantic_chunking=enable_semantic_chunking
            )
            if doc:
                documents.append(doc)
                chunk_info = f" + {doc['chunk_count']} chunks" if doc.get('chunks') else ""
                logger.info(f"✓ Processed: {doc_file.name} ({doc['word_count']} words{chunk_info})")
            else:
                logger.warning(f"✗ Failed: {doc_file.name}")
        
        logger.info(f"Successfully processed {len(documents)}/{len(all_files)} documents")
        return documents


# Backward compatibility alias
PDFProcessor = DocumentProcessor
